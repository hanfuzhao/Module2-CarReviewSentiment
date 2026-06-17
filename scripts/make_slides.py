import json
import os
import re
import zipfile
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path("PITCH.pptx")
METRICS = json.loads(Path("data/outputs/metrics.json").read_text())
PLOTS = Path("data/outputs/plots")
BG = RGBColor(15, 17, 21)
PANEL = RGBColor(26, 29, 36)
TEXT = RGBColor(232, 234, 237)
MUTED = RGBColor(154, 160, 170)
ACCENT = RGBColor(74, 158, 255)
NEG = RGBColor(217, 83, 79)
NEU = RGBColor(240, 173, 78)
POS = RGBColor(92, 184, 92)
WHITE = RGBColor(255, 255, 255)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def pct(x):
    return "n/a" if x is None else f"{x * 100:.0f}%"


def pct1(x):
    return "n/a" if x is None else f"{x * 100:.1f}%"


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def box(s, left, top, width, height):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def para(
    tf,
    text,
    size=18,
    color=TEXT,
    bold=False,
    first=False,
    align=PP_ALIGN.LEFT,
    space_after=6,
    bullet=False,
):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    runs = text if isinstance(text, list) else [(text, color, bold)]
    for i, item in enumerate(runs):
        t, c, b = item if isinstance(item, tuple) else (item, color, bold)
        r = p.add_run()
        r.text = ("- " if bullet and i == 0 else "") + t
        r.font.size = Pt(size)
        r.font.color.rgb = c
        r.font.bold = b
        r.font.name = "Helvetica Neue"
    return p


def title(s, text, accent_tail=None):
    tf = box(s, 0.6, 0.35, 12.1, 1.0)
    runs = [(text, TEXT, True)]
    if accent_tail:
        runs.append((accent_tail, ACCENT, True))
    para(tf, runs, size=30, first=True)
    ln = s.shapes.add_shape(1, Inches(0.65), Inches(1.25), Inches(1.6), Pt(3))
    ln.fill.solid()
    ln.fill.fore_color.rgb = ACCENT
    ln.line.fill.background()


def chip(s, left, top, text, color):
    w = 0.18 + 0.105 * len(text)
    sh = s.shapes.add_shape(5, Inches(left), Inches(top), Inches(w), Inches(0.42))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = WHITE
    return left + w + 0.15


s = slide()
tf = box(s, 0.8, 2.3, 11.7, 2.2)
para(
    tf,
    [("Car Review Sentiment Analyzer", TEXT, True)],
    size=44,
    first=True,
)
para(
    tf,
    "Per-aspect sentiment for owner reviews, with evaluation at the core.",
    size=20,
    color=MUTED,
    space_after=18,
)
para(
    tf,
    [
        ("Can Machines Understand Us Reliably?", MUTED, False),
        ("     Module 2 Hackathon", MUTED, False),
    ],
    size=16,
)
tf2 = box(s, 0.8, 5.5, 11.7, 1.5)
para(tf2, "Hanfu Zhao", size=18, color=TEXT, bold=True, first=True)
para(tf2, "Live demo:  https://HanfuZhao781-car-review-sentiment.hf.space", size=14, color=ACCENT)
para(tf2, "Code:  https://github.com/hanfuzhao/Module2-CarReviewSentiment", size=14, color=MUTED)
s = slide()
title(s, "The problem")
tf = box(s, 0.65, 1.7, 12.0, 5.2)
para(
    tf,
    [
        ("Platforms like ", TEXT, False),
        ("Dongchedi", ACCENT, True),
        (" turn thousands of owner reviews into ", TEXT, False),
        ("per-dimension", ACCENT, True),
        (" ratings for engine, comfort, and value that buyers trust.", TEXT, False),
    ],
    size=20,
    first=True,
    space_after=14,
)
para(
    tf,
    "Those ratings come from a sentiment model. Misread one review and the error is averaged into a public rating thousands of people see.",
    size=20,
    color=TEXT,
    space_after=18,
)
para(tf, "Two notorious failure modes:", size=20, bold=True, color=TEXT, space_after=10)
para(
    tf,
    [("Negation", NEG, True), (", the brakes are not good", TEXT, False)],
    size=19,
    bullet=True,
    space_after=8,
)
para(
    tf,
    [("Sarcasm", NEG, True), (", great, another trip to the dealer this week", TEXT, False)],
    size=19,
    bullet=True,
    space_after=18,
)
para(
    tf,
    [
        ("So this is a sentiment classifier where the real work is ", MUTED, False),
        ("measuring where it breaks.", ACCENT, True),
    ],
    size=20,
)
s = slide()
title(s, "Approach: ", "3 models plus transfer learning")
tf = box(s, 0.65, 1.7, 12.0, 1.2)
para(
    tf,
    [
        ("Data:  ", MUTED, True),
        (
            "36,518 real Edmunds owner reviews. 3-class labels from the 1 to 5 star rating.",
            TEXT,
            False,
        ),
    ],
    size=18,
    first=True,
)
cards = [
    ("Naive", "majority class\nthe floor", PANEL),
    ("Classical", "TF-IDF plus logistic\nregression", PANEL),
    ("DistilBERT", "fine-tuned\nDEPLOYED", RGBColor(30, 51, 74)),
]
x = 0.65
for name, desc, col in cards:
    card = s.shapes.add_shape(5, Inches(x), Inches(3.05), Inches(3.7), Inches(1.9))
    card.fill.solid()
    card.fill.fore_color.rgb = col
    card.line.color.rgb = ACCENT if "DEPLOY" in desc else PANEL
    ctf = card.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(
        ctf,
        name,
        size=22,
        bold=True,
        color=ACCENT if "DEPLOY" in desc else TEXT,
        first=True,
        align=PP_ALIGN.CENTER,
        space_after=6,
    )
    for line in desc.split("\n"):
        para(ctf, line, size=15, color=MUTED, align=PP_ALIGN.CENTER, space_after=2)
    x += 3.95
tf = box(s, 0.65, 5.3, 12.0, 1.7)
para(
    tf,
    [
        ("Transfer learning:  ", ACCENT, True),
        (
            "start from distilbert-base-uncased, pretrained on general English, add a sentiment head, fine-tune on car reviews.",
            TEXT,
            False,
        ),
    ],
    size=18,
    first=True,
    space_after=10,
)
para(
    tf,
    [
        ("Plus a per-aspect ", MUTED, False),
        ("breakdown", ACCENT, True),
        (", lexicon-routed clauses scored by the same model.", MUTED, False),
    ],
    size=16,
)
s = slide()
title(s, "Data augmentation")
tf = box(s, 0.65, 1.7, 12.0, 5.2)
para(
    tf,
    [
        ("The catch:  ", NEU, True),
        (
            "the data is 86% positive. A model that always says positive scores about 86% accuracy and is useless.",
            TEXT,
            False,
        ),
    ],
    size=20,
    first=True,
    space_after=18,
)
para(
    tf,
    "To rescue the starved minority classes, EDA, Easy Data Augmentation:",
    size=19,
    color=TEXT,
    space_after=10,
)
para(tf, "synonym replacement with WordNet", size=18, bullet=True, space_after=6)
para(tf, "random word swap and random deletion", size=18, bullet=True, space_after=6)
para(
    tf,
    "back-translation, English to German to English, for paraphrase diversity",
    size=18,
    bullet=True,
    space_after=18,
)
para(
    tf,
    [
        ("The interesting question is not whether to augment, it is ", MUTED, False),
        ("when it helps. Evaluation answered that.", ACCENT, True),
    ],
    size=20,
)
m = METRICS["models"]
s = slide()
title(s, "Results: ", "accuracy lies")
rows = [("Model", "Accuracy", "Macro-F1", "Neg. recall")]
for k, label in (
    ("naive", "Naive, majority"),
    ("classical", "Classical TF-IDF"),
    ("deep", "DistilBERT, deployed"),
):
    d = m[k]
    rows.append(
        (label, pct(d["accuracy"]), pct1(d["macro_f1"]), pct1(d["per_class"]["negative"]["recall"]))
    )
tbl_shape = s.shapes.add_table(len(rows), 4, Inches(0.65), Inches(1.7), Inches(7.4), Inches(2.6))
table = tbl_shape.table
table.columns[0].width = Inches(3.1)
for j in range(1, 4):
    table.columns[j].width = Inches(1.43)
for i, row in enumerate(rows):
    for j, val in enumerate(row):
        cell = table.cell(i, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = (
            RGBColor(35, 51, 69) if i == 0 else RGBColor(30, 51, 74) if i == 3 else PANEL
        )
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        r = p.add_run()
        r.text = val
        r.font.size = Pt(15)
        r.font.color.rgb = TEXT
        r.font.bold = i == 0 or i == 3
tf = box(s, 0.65, 4.6, 7.5, 2.6)
para(
    tf,
    [
        ("Naive: 86% accuracy, ", TEXT, False),
        ("0% negative recall.", NEG, True),
        (" Accuracy is the wrong metric on imbalanced data.", TEXT, False),
    ],
    size=17,
    first=True,
    space_after=12,
)
para(
    tf,
    [("What matters is macro-F1, which shows transfer learning working:", TEXT, False)],
    size=17,
    space_after=8,
)
para(
    tf,
    [
        (f"{pct1(m['naive']['macro_f1'])}", MUTED, True),
        ("  to  ", MUTED, False),
        (f"{pct1(m['classical']['macro_f1'])}", TEXT, True),
        ("  to  ", MUTED, False),
        (f"{pct1(m['deep']['macro_f1'])}", ACCENT, True),
        ("   naive to classical to DistilBERT", MUTED, False),
    ],
    size=18,
)
img = PLOTS / "class_distribution.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(8.4), Inches(1.9), height=Inches(3.2))
    tf = box(s, 8.4, 5.1, 4.4, 0.6)
    para(tf, "14 to 1 class imbalance", size=13, color=MUTED, first=True, align=PP_ALIGN.CENTER)
ae = METRICS["augmentation_effect"]
aed = METRICS["augmentation_effect_deep"]
s = slide()
title(s, "When does augmentation help?")
tf = box(s, 0.65, 1.55, 12.0, 0.9)
para(
    tf,
    [
        ("Cold-start regime, a newly-launched car with only ", TEXT, False),
        (f"{METRICS['config']['lowres_cap']:,} reviews", ACCENT, True),
        (", realistic for Dongchedi:", TEXT, False),
    ],
    size=18,
    first=True,
)


def panel(left, head, head_color, before, after, verdict, verdict_color):
    card = s.shapes.add_shape(5, Inches(left), Inches(2.6), Inches(5.7), Inches(3.4))
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL
    card.line.color.rgb = head_color
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.25)
    para(tf, head, size=20, bold=True, color=head_color, first=True, space_after=14)
    para(
        tf,
        [
            ("macro-F1:   ", MUTED, False),
            (before, TEXT, True),
            ("  to  ", MUTED, False),
            (after, head_color, True),
        ],
        size=20,
        space_after=10,
    )
    para(tf, verdict, size=18, color=verdict_color, bold=True)
    return card


panel(
    0.65,
    "Classical model",
    POS,
    pct1(ae["macro_f1"]["before"]),
    pct1(ae["macro_f1"]["after"]),
    "Augmentation wins",
    POS,
)
panel(
    6.95,
    "Transformer",
    MUTED,
    pct1(aed["macro_f1"]["before"]),
    pct1(aed["macro_f1"]["after"]),
    "Already strong, no help",
    NEU,
)
tf = box(s, 0.65, 6.25, 12.0, 1.0)
para(
    tf,
    [
        ("Same technique, opposite verdict. ", TEXT, False),
        ("Only rigorous evaluation could tell them apart.", ACCENT, True),
    ],
    size=19,
    first=True,
    align=PP_ALIGN.CENTER,
)
neg = METRICS["stress_tests"]["negation"]["deep"]
sar = METRICS["stress_tests"]["sarcasm"]["deep"]
gate = METRICS["confidence_gating"]
g0 = gate[0]["accuracy"]
gbest = max(
    (c for c in gate if c["accuracy"] is not None and c["coverage"] >= 0.5),
    key=lambda c: c["accuracy"],
)
s = slide()
title(s, "Stress tests ", "reveal the hard limits")
tf = box(s, 0.65, 1.65, 7.4, 5.2)
para(
    tf,
    "Hand-built linguistic stress tests expose what aggregate accuracy hides:",
    size=18,
    color=TEXT,
    first=True,
    space_after=16,
)
para(
    tf,
    [
        ("Negation:  ", TEXT, True),
        ("affirmative ", TEXT, False),
        (pct(neg["affirmative_accuracy"]), POS, True),
        ("  versus negated ", TEXT, False),
        (pct(neg["negated_accuracy"]), NEG, True),
    ],
    size=19,
    bullet=True,
    space_after=12,
)
para(
    tf,
    [
        ("Sarcasm:  ", TEXT, True),
        (pct(sar["misread_as_positive_rate"]), NEG, True),
        (" misread as positive", TEXT, False),
    ],
    size=19,
    bullet=True,
    space_after=18,
)
para(
    tf,
    [("The deployable fix, confidence-gated abstention:", ACCENT, True)],
    size=19,
    bold=True,
    space_after=10,
)
para(
    tf,
    [
        ("accuracy ", TEXT, False),
        (pct1(g0), TEXT, True),
        ("  to  ", MUTED, False),
        (pct1(gbest["accuracy"]), POS, True),
        (f"  at {pct(gbest['coverage'])} coverage", MUTED, False),
    ],
    size=19,
    space_after=8,
)
para(tf, "low-confidence reviews are flagged for a human.", size=16, color=MUTED)
for img, top in ((PLOTS / "confusion_deep.png", 1.75), (PLOTS / "abstention_curve.png", 4.4)):
    if img.exists():
        s.shapes.add_picture(str(img), Inches(8.55), Inches(top), height=Inches(2.5))
s = slide()
tf = box(s, 0.8, 1.6, 11.7, 1.4)
para(tf, [("Live demo", ACCENT, True)], size=34, first=True)
tf = box(s, 0.8, 2.9, 11.7, 3.0)
para(
    tf,
    "Paste a review to get overall sentiment plus the per-aspect breakdown, all from one review.",
    size=20,
    color=TEXT,
    first=True,
    space_after=12,
)
para(
    tf,
    "Mixed review: engine positive, fuel economy negative, price negative, comfort positive",
    size=17,
    color=MUTED,
    space_after=6,
    bullet=True,
)
para(
    tf,
    "Not bad: deep model correct, classical model wrong, the model gap shown live",
    size=17,
    color=MUTED,
    space_after=6,
    bullet=True,
)
para(
    tf,
    "Heavy negation: the app abstains. Sarcasm: confidently wrong, the honest frontier",
    size=17,
    color=MUTED,
    space_after=18,
    bullet=True,
)
para(
    tf,
    [
        ("Try it live:  ", TEXT, True),
        ("https://HanfuZhao781-car-review-sentiment.hf.space", ACCENT, True),
    ],
    size=20,
)
tf = box(s, 0.8, 6.1, 11.7, 1.0)
para(
    tf,
    [("Accurate, honest, and deployable. ", TEXT, True), ("Thank you.", ACCENT, True)],
    size=22,
    first=True,
)
n_slides = len(prs.slides._sldIdLst)
prs.save(str(OUT))


def strip_non_ascii_theme(path):
    tmp = str(path) + ".tmp"
    with zipfile.ZipFile(path) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.namelist():
            data = zin.read(item)
            if item.endswith(".xml"):
                txt = data.decode("utf-8")
                txt = re.sub(r'<a:font script="[^"]*" typeface="[^"]*"/>', "", txt)
                txt = re.sub(r'(typeface=")[^"]*"', lambda m: m.group(1) + '"'
                             if any(ord(c) > 127 for c in m.group(0)) else m.group(0), txt)
                data = txt.encode("utf-8")
            zout.writestr(item, data)
    os.replace(tmp, path)


strip_non_ascii_theme(OUT)
print(f"Wrote {OUT}  ({n_slides} slides)")
